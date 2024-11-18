import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
import io
import tempfile
import os
import logging
from wordcloud import WordCloud
import matplotlib.pyplot as plt
import base64
import numpy as np
from collections import Counter
import pandas as pd
from datetime import datetime

# Set up logging for better debugging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def generate_word_cloud(df, text_column, title):
    """Generate a word cloud from text data"""
    # Combine all text into one string
    text = ' '.join(df[text_column].astype(str).fillna(''))
    
    # Generate word cloud
    wordcloud = WordCloud(
        width=800, 
        height=400,
        background_color='white',
        colormap='viridis',
        max_words=100
    ).generate(text)
    
    # Create plotly figure
    img_bytes = io.BytesIO()
    plt.figure(figsize=(10, 5))
    plt.imshow(wordcloud, interpolation='bilinear')
    plt.axis('off')
    plt.title(title)
    plt.savefig(img_bytes, format='png', bbox_inches='tight', pad_inches=0)
    plt.close()
    
    # Convert to base64
    img_bytes.seek(0)
    img_base64 = base64.b64encode(img_bytes.read()).decode()
    
    # Create plotly figure with image
    fig = go.Figure()
    fig.add_layout_image(
        dict(
            source=f'data:image/png;base64,{img_base64}',
            x=0,
            y=1,
            sizex=1,
            sizey=1,
            sizing="stretch",
            layer="below"
        )
    )
    
    # Update layout
    fig.update_layout(
        title=title,
        showlegend=False,
        width=800,
        height=400,
        margin=dict(l=0, r=0, t=30, b=0)
    )
    
    # Remove axes
    fig.update_xaxes(showticklabels=False, showgrid=False, zeroline=False)
    fig.update_yaxes(showticklabels=False, showgrid=False, zeroline=False)
    
    return fig

def generate_chart(df, chart_type, x_column=None, y_column=None, title=None, text_column=None):
    """Generate various types of charts based on the specified type"""
    if chart_type == "word_cloud":
        return generate_word_cloud(df, text_column, title)
        
    try:
        if chart_type == "line":
            fig = px.line(df, x=x_column, y=y_column, title=title)
        elif chart_type == "bar":
            fig = px.bar(df, x=x_column, y=y_column, title=title)
        elif chart_type == "scatter":
            fig = px.scatter(df, x=x_column, y=y_column, title=title)
        elif chart_type == "pie":
            fig = px.pie(df, names=x_column, values=y_column, title=title)
        else:
            return None
        
        # Improve chart appearance
        fig.update_layout(
            title_x=0.5,
            margin=dict(t=50, l=50, r=50, b=50),
            template="plotly_white",
            height=500,
            width=800,
            showlegend=True
        )
        
        if chart_type != "pie":
            fig.update_xaxes(title_text=x_column.replace('_', ' ').title())
            fig.update_yaxes(title_text=y_column.replace('_', ' ').title())
        
        return fig
    except Exception as e:
        logger.error(f"Error generating chart: {str(e)}")
        return None

def save_chart_as_image(chart):
    """Save Plotly chart as image and return the path"""
    try:
        # Create a temporary file with a unique name
        temp_dir = tempfile.gettempdir()
        temp_path = os.path.join(temp_dir, f'chart_{datetime.now().strftime("%Y%m%d_%H%M%S_%f")}.png')
        
        logger.info(f"Attempting to save chart as image to: {temp_path}")
        
        # Save the chart as a static image
        chart.write_image(
            temp_path,
            format="png",
            engine="kaleido",
            width=1600,  # Increased resolution
            height=900,   # Increased resolution
            scale=2
        )
        
        # Verify the file exists and has content
        if os.path.exists(temp_path) and os.path.getsize(temp_path) > 0:
            logger.info(f"Chart successfully saved to {temp_path}")
            return temp_path
        else:
            logger.error("Chart file was not created or is empty")
            return None
            
    except Exception as e:
        logger.error(f"Error saving chart as image: {str(e)}")
        return None

def add_text_to_slide(slide, text, is_title=False):
    """Add text to a slide with proper formatting"""
    if is_title:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        font_size = Pt(32)
    else:
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(4)
        font_size = Pt(14)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = text
    
    paragraph = tf.paragraphs[0]
    paragraph.font.size = font_size
    if is_title:
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(44, 62, 80)
    
    tf.word_wrap = True

def create_presentation(messages):
    """Create a PowerPoint presentation from chat messages"""
    prs = Presentation()
    temp_files = []  # Keep track of temporary files
    
    try:
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Data Analysis Report"
        if hasattr(title_slide.shapes, 'subtitle') and title_slide.shapes.subtitle:
            title_slide.shapes.subtitle.text = "Generated by AI Assistant"
        
        for message in messages:
            if message["role"] == "assistant":
                # Create a new slide for text content
                text_slide = prs.slides.add_slide(prs.slide_layouts[6])
                add_text_to_slide(text_slide, "Analysis", is_title=True)
                add_text_to_slide(text_slide, str(message["content"]))
                
                # Handle charts
                if "chart" in message and message["chart"] is not None:
                    charts_to_process = []
                    
                    # Handle both single charts and lists of charts
                    if isinstance(message["chart"], list):
                        charts_to_process.extend(message["chart"])
                    else:
                        charts_to_process.append(message["chart"])
                    
                    # Process each chart
                    for chart in charts_to_process:
                        try:
                            logger.info("Processing chart for PowerPoint...")
                            
                            # Save the chart as an image
                            chart_image_path = save_chart_as_image(chart)
                            if chart_image_path:
                                temp_files.append(chart_image_path)
                                
                                # Create chart slide
                                chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
                                
                                # Try to get chart title from layout
                                chart_title = "Visualization"
                                try:
                                    if hasattr(chart, 'layout') and chart.layout.title:
                                        chart_title = chart.layout.title.text
                                except:
                                    pass
                                
                                add_text_to_slide(chart_slide, chart_title, is_title=True)
                                
                                # Add the chart image
                                left = Inches(1)
                                top = Inches(1.5)
                                width = Inches(8)
                                chart_slide.shapes.add_picture(chart_image_path, left, top, width=width)
                                logger.info(f"Chart '{chart_title}' successfully added to slide")
                            else:
                                logger.error("Failed to save chart as image")
                                
                        except Exception as e:
                            logger.error(f"Error adding chart to slide: {str(e)}")
                            continue
        
        # Save the presentation
        temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        prs.save(temp_pptx.name)
        logger.info(f"Presentation saved successfully to {temp_pptx.name}")
        
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
                logger.info(f"Cleaned up temporary file: {temp_file}")
            except Exception as e:
                logger.error(f"Error cleaning up temporary file {temp_file}: {str(e)}")
        
        return temp_pptx.name
        
    except Exception as e:
        logger.error(f"Error creating presentation: {str(e)}")
        
        # Clean up temporary files in case of error
        for temp_file in temp_files:
            try:
                os.remove(temp_file)
            except:
                pass
                
        return None